#!/usr/bin/env python3
"""Extract text from AIDD PDF and PPTX source documents."""
from __future__ import annotations

import json
from pathlib import Path

import fitz
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]


def extract_pdf(pdf_path: Path, output_path: Path) -> None:
    doc = fitz.open(pdf_path)
    pages = []
    for i, page in enumerate(doc, 1):
        text = page.get_text("text").strip()
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        pages.append({"page": i, "text": text, "preview": " | ".join(lines[:20])})
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(json.dumps(pages, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"PDF pages {len(pages)}, nonempty {sum(bool(p['text']) for p in pages)} -> {output_path}")


def extract_pptx(pptx_path: Path, output_path: Path) -> None:
    prs = Presentation(str(pptx_path))
    slides = []
    for si, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.extend(line.strip() for line in shape.text.splitlines() if line.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text:
                            texts.extend(line.strip() for line in cell.text.splitlines() if line.strip())
        slides.append({"slide": si, "text": "\n".join(texts), "preview": " | ".join(texts[:20])})
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(json.dumps(slides, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"PPT slides {len(slides)}, nonempty {sum(bool(s['text']) for s in slides)} -> {output_path}")


def main() -> None:
    extract_pdf(
        ROOT / "aidd分享图片/网易/AiDD2026上海站-赵雨森-052302.pdf",
        ROOT / "output/aidd_网易_pdf_text.json",
    )
    extract_pptx(
        ROOT / "aidd分享图片/中兴通讯2/丁晓彬-AiDD2026上海站-v1.pptx",
        ROOT / "output/aidd_中兴通讯2_pptx_text.json",
    )


if __name__ == "__main__":
    main()
